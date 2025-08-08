"""
Enhanced Document Processing Service
Phase 9: Comprehensive Knowledge File Support

Supports:
- Microsoft Office: Word (.docx), Excel (.xlsx/.xls), PowerPoint (.pptx)
- PDF documents
- Google Workspace: Docs, Sheets, Slides
- Rich Text Format (.rtf)
- Various text formats: .txt, .log, .xml, .yaml, .ini, .properties, .sql
- Source code files: .py, .js, .html, .css, .java, .cpp, .php, .rb, .go, .rs, .ts
"""

import logging
import io
import re
import json
import csv
from typing import Dict, Any, Optional, List, Union
from pathlib import Path
import chardet

# Microsoft Office processing
try:
    from docx import Document as DocxDocument  # type: ignore
    import openpyxl  # type: ignore
    from pptx import Presentation  # type: ignore
    import xlrd  # type: ignore
except ImportError:
    DocxDocument = None  # type: ignore
    openpyxl = None  # type: ignore
    Presentation = None  # type: ignore
    xlrd = None  # type: ignore

# PDF processing
try:
    import PyPDF2  # type: ignore
    import pdfplumber  # type: ignore
except ImportError:
    PyPDF2 = None  # type: ignore
    pdfplumber = None  # type: ignore

# Rich Text and other formats
try:
    from striprtf.striprtf import rtf_to_text  # type: ignore
    import yaml  # type: ignore
    from bs4 import BeautifulSoup  # type: ignore
    import configparser
except ImportError:
    rtf_to_text = None  # type: ignore
    yaml = None  # type: ignore
    BeautifulSoup = None  # type: ignore
    configparser = None  # type: ignore

# Google API
try:
    from google.oauth2.credentials import Credentials  # type: ignore
    from googleapiclient.discovery import build  # type: ignore
    from google.auth.transport.requests import Request  # type: ignore
    from google_auth_oauthlib.flow import InstalledAppFlow  # type: ignore
except ImportError:
    Credentials = None  # type: ignore
    build = None  # type: ignore
    Request = None  # type: ignore
    InstalledAppFlow = None  # type: ignore

logger = logging.getLogger(__name__)

class DocumentProcessorError(Exception):
    """Document processing related errors"""
    pass

class DocumentProcessor:
    """Enhanced document processor with chunking capabilities"""
    
    def __init__(self):
        """Initialize document processor"""
        self.supported_formats = {
            '.docx': self._process_docx,
            '.pdf': self._process_pdf,
            '.xlsx': self._process_xlsx,
            '.xls': self._process_xls,
            '.pptx': self._process_pptx,
            '.rtf': self._process_rtf,
            '.txt': self._process_text,
            '.log': self._process_text,
            '.md': self._process_markdown,
            '.markdown': self._process_markdown,
            '.json': self._process_json,
            '.yaml': self._process_yaml,
            '.yml': self._process_yaml,
            '.xml': self._process_xml,
            '.html': self._process_html,
            '.htm': self._process_html,
            '.css': self._process_text,
            '.js': self._process_text,
            '.py': self._process_text,
            '.java': self._process_text,
            '.cpp': self._process_text,
            '.c': self._process_text,
            '.h': self._process_text,
            '.cs': self._process_text,
            '.php': self._process_text,
            '.rb': self._process_text,
            '.go': self._process_text,
            '.rs': self._process_text,
            '.ts': self._process_text,
            '.sh': self._process_text,
            '.sql': self._process_text,
            '.ini': self._process_ini,
            '.properties': self._process_properties,
            '.csv': self._process_csv,
            '.tsv': self._process_csv,
        }
        logger.info("Document processor initialized")
        
        # チャンク化設定
        self.chunk_size = 4000  # 4,000文字単位
        self.chunk_overlap = 200  # 200文字のオーバーラップ
        self.max_chunks_per_file = 50  # 1ファイル最大50チャンク
    
    async def process_document(
        self, 
        file_content: bytes, 
        filename: str, 
        content_type: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Process document and extract text content
        
        Args:
            file_content: Raw file content as bytes
            filename: Original filename
            content_type: MIME content type
            
        Returns:
            Processed document information
        """
        try:
            file_path = Path(filename)
            file_extension = file_path.suffix.lower()
            
            if file_extension not in self.supported_formats:
                raise DocumentProcessorError(f"Unsupported file format: {file_extension}")
            
            # Process the document
            processor_func = self.supported_formats[file_extension]
            result = await processor_func(file_content, filename)
            
            # Add metadata
            result.update({
                "filename": filename,
                "file_extension": file_extension,
                "content_type": content_type,
                "file_size": len(file_content),
                "processing_status": "success"
            })
            
            logger.info(f"Successfully processed {filename} ({file_extension})")
            return result
            
        except Exception as e:
            logger.error(f"Failed to process document {filename}: {str(e)}")
            raise DocumentProcessorError(f"Document processing failed: {str(e)}")
    
    def _detect_encoding(self, content: bytes) -> str:
        """Detect text encoding"""
        try:
            detection = chardet.detect(content)
            encoding = detection.get('encoding', 'utf-8') if detection else 'utf-8'
            return encoding or 'utf-8'
        except Exception:
            return 'utf-8'
    
    async def _process_docx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Word document (.docx)"""
        if DocxDocument is None:
            raise DocumentProcessorError("python-docx not installed")
        
        try:
            doc = DocxDocument(io.BytesIO(content))
            
            # Extract text from paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            tables_content = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        tables_content.append(" | ".join(row_text))
            
            full_text = "\n".join(text_content)
            if tables_content:
                full_text += "\n\nTables:\n" + "\n".join(tables_content)
            
            return {
                "text": full_text,
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
                "word_count": len(full_text.split()),
                "document_type": "word_document"
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process Word document: {str(e)}")
    
    async def _process_pdf(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PDF document"""
        text_content = ""
        page_count = 0
        
        # Try pdfplumber first (better for complex layouts)
        if pdfplumber:
            try:
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    page_count = len(pdf.pages)
                    text_parts = []
                    
                    for page in pdf.pages:
                        page_text = page.extract_text()  # type: ignore
                        if page_text:
                            text_parts.append(page_text)
                    
                    text_content = "\n\n".join(text_parts)
                    
                    return {
                        "text": text_content,
                        "pages": page_count,
                        "word_count": len(text_content.split()),
                        "document_type": "pdf_document",
                        "extraction_method": "pdfplumber"
                    }
            except Exception as e:
                logger.warning(f"pdfplumber failed, trying PyPDF2: {str(e)}")
        
        # Fallback to PyPDF2
        if PyPDF2:
            try:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                page_count = len(pdf_reader.pages)
                text_parts = []
                
                for page in pdf_reader.pages:  # type: ignore
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                
                text_content = "\n\n".join(text_parts)
                
                return {
                    "text": text_content,
                    "pages": page_count,
                    "word_count": len(text_content.split()),
                    "document_type": "pdf_document",
                    "extraction_method": "PyPDF2"
                }
            except Exception as e:
                raise DocumentProcessorError(f"Failed to process PDF: {str(e)}")
        
        raise DocumentProcessorError("No PDF processing library available")
    
    async def _process_xlsx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Excel file (.xlsx)"""
        if openpyxl is None:
            raise DocumentProcessorError("openpyxl not installed")
        
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
            sheets_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):  # Skip empty rows
                        sheet_data.append(" | ".join(row_data))
                
                if sheet_data:
                    sheets_content.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_data))
            
            full_text = "\n\n".join(sheets_content)
            
            return {
                "text": full_text,
                "sheets": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames,
                "word_count": len(full_text.split()),
                "document_type": "excel_document"
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process Excel file: {str(e)}")
    
    async def _process_xls(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process legacy Excel file (.xls)"""
        if xlrd is None:
            raise DocumentProcessorError("xlrd not installed")
        
        try:
            workbook = xlrd.open_workbook(file_contents=content)
            sheets_content = []
            
            for sheet_index in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_index)
                sheet_name = workbook.sheet_names()[sheet_index]
                sheet_data = []
                
                for row_idx in range(sheet.nrows):
                    row_data = []
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        row_data.append(str(cell_value) if cell_value else "")
                    
                    if any(cell.strip() for cell in row_data):
                        sheet_data.append(" | ".join(row_data))
                
                if sheet_data:
                    sheets_content.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_data))
            
            full_text = "\n\n".join(sheets_content)
            
            return {
                "text": full_text,
                "sheets": workbook.nsheets,
                "sheet_names": workbook.sheet_names(),
                "word_count": len(full_text.split()),
                "document_type": "excel_legacy_document"
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process legacy Excel file: {str(e)}")
    
    async def _process_pptx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PowerPoint presentation (.pptx)"""
        if Presentation is None:
            raise DocumentProcessorError("python-pptx not installed")
        
        try:
            prs = Presentation(io.BytesIO(content))
            slides_content = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    slides_content.append(f"Slide {slide_idx}:\n" + "\n".join(slide_text))
            
            full_text = "\n\n".join(slides_content)
            
            return {
                "text": full_text,
                "slides": len(prs.slides),
                "word_count": len(full_text.split()),
                "document_type": "powerpoint_document"
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process PowerPoint file: {str(e)}")
    
    async def _process_rtf(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Rich Text Format (.rtf)"""
        if rtf_to_text is None:
            raise DocumentProcessorError("striprtf not installed")
        
        try:
            encoding = self._detect_encoding(content)
            rtf_content = content.decode(encoding, errors='replace')
            text_content = rtf_to_text(rtf_content)
            
            return {
                "text": text_content,
                "word_count": len(text_content.split()),
                "document_type": "rtf_document",
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process RTF file: {str(e)}")
    
    async def _process_text(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process plain text files"""
        try:
            encoding = self._detect_encoding(content)
            text_content = content.decode(encoding, errors='replace')
            
            # Detect file type based on extension
            file_path = Path(filename)
            extension = file_path.suffix.lower()
            
            document_types = {
                '.txt': 'text_document',
                '.log': 'log_file',
                '.py': 'python_source',
                '.js': 'javascript_source',
                '.java': 'java_source',
                '.cpp': 'cpp_source',
                '.c': 'c_source',
                '.h': 'header_file',
                '.cs': 'csharp_source',
                '.php': 'php_source',
                '.rb': 'ruby_source',
                '.go': 'go_source',
                '.rs': 'rust_source',
                '.ts': 'typescript_source',
                '.sh': 'shell_script',
                '.sql': 'sql_script',
                '.css': 'css_stylesheet'
            }
            
            return {
                "text": text_content,
                "lines": len(text_content.splitlines()),
                "word_count": len(text_content.split()),
                "character_count": len(text_content),
                "document_type": document_types.get(extension, 'text_document'),
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process text file: {str(e)}")
    
    async def _process_markdown(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Markdown files"""
        try:
            encoding = self._detect_encoding(content)
            text_content = content.decode(encoding, errors='replace')
            
            # Remove markdown formatting for analysis
            clean_text = re.sub(r'[#*`\[\]()_~]', '', text_content)
            clean_text = re.sub(r'!\[.*?\]\(.*?\)', '', clean_text)  # Remove images
            clean_text = re.sub(r'\[.*?\]\(.*?\)', '', clean_text)   # Remove links
            
            return {
                "text": text_content,
                "clean_text": clean_text,
                "lines": len(text_content.splitlines()),
                "word_count": len(clean_text.split()),
                "document_type": "markdown_document",
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process Markdown file: {str(e)}")
    
    async def _process_json(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process JSON files"""
        try:
            encoding = self._detect_encoding(content)
            json_content = content.decode(encoding, errors='replace')
            
            # Parse JSON and convert to readable text
            data = json.loads(json_content)
            
            # Convert JSON to readable format
            readable_text = json.dumps(data, indent=2, ensure_ascii=False)
            
            return {
                "text": json_content,
                "readable_text": readable_text,
                "json_structure": str(type(data).__name__),
                "word_count": len(readable_text.split()),
                "document_type": "json_document",
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process JSON file: {str(e)}")
    
    async def _process_yaml(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process YAML files"""
        if yaml is None:
            return await self._process_text(content, filename)
        
        try:
            encoding = self._detect_encoding(content)
            yaml_content = content.decode(encoding, errors='replace')
            
            # Parse YAML and convert to readable text
            data = yaml.safe_load(yaml_content)
            readable_text = yaml.dump(data, default_flow_style=False, allow_unicode=True)
            
            return {
                "text": yaml_content,
                "readable_text": readable_text,
                "yaml_structure": str(type(data).__name__),
                "word_count": len(readable_text.split()),
                "document_type": "yaml_document",
                "encoding": encoding
            }
            
        except Exception as e:
            # Fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_xml(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process XML files"""
        try:
            encoding = self._detect_encoding(content)
            xml_content = content.decode(encoding, errors='replace')
            
            # Parse XML and extract text
            if BeautifulSoup is not None:
                try:
                    soup = BeautifulSoup(xml_content, 'xml')
                    text_content = soup.get_text(separator=' ', strip=True)
                    
                    return {
                        "text": xml_content,
                        "extracted_text": text_content,
                        "word_count": len(text_content.split()),
                        "document_type": "xml_document",
                        "encoding": encoding
                    }
                except Exception as e:
                    logger.warning(f"BeautifulSoup XML parsing failed for {filename}: {str(e)}")
            
            # Fallback: Simple XML tag removal
            import re
            # Remove XML tags
            clean_text = re.sub(r'<[^>]+>', ' ', xml_content)
            # Clean up whitespace
            clean_text = re.sub(r'\s+', ' ', clean_text).strip()
            
            return {
                "text": xml_content,
                "extracted_text": clean_text,
                "word_count": len(clean_text.split()),
                "document_type": "xml_document",
                "encoding": encoding
            }
            
        except Exception as e:
            logger.error(f"XML processing failed for {filename}: {str(e)}")
            # Final fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_html(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process HTML files"""
        try:
            encoding = self._detect_encoding(content)
            html_content = content.decode(encoding, errors='replace')
            
            # Parse HTML and extract text
            if BeautifulSoup is not None:
                try:
                    soup = BeautifulSoup(html_content, 'html.parser')
                    
                    # Remove script and style elements
                    for script in soup(["script", "style"]):
                        script.decompose()
                    
                    text_content = soup.get_text(separator=' ', strip=True)
                    title = soup.title.string if soup.title else None
                    
                    return {
                        "text": html_content,
                        "extracted_text": text_content,
                        "title": title,
                        "word_count": len(text_content.split()),
                        "document_type": "html_document",
                        "encoding": encoding
                    }
                except Exception as e:
                    logger.warning(f"BeautifulSoup HTML parsing failed for {filename}: {str(e)}")
            
            # Fallback: Simple HTML tag removal
            import re
            # Remove HTML tags
            clean_text = re.sub(r'<[^>]+>', ' ', html_content)
            # Clean up whitespace
            clean_text = re.sub(r'\s+', ' ', clean_text).strip()
            
            return {
                "text": html_content,
                "extracted_text": clean_text,
                "title": None,
                "word_count": len(clean_text.split()),
                "document_type": "html_document",
                "encoding": encoding
            }
            
        except Exception as e:
            logger.error(f"HTML processing failed for {filename}: {str(e)}")
            # Final fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_ini(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process INI/config files"""
        if configparser is None:
            return await self._process_text(content, filename)
        
        try:
            encoding = self._detect_encoding(content)
            ini_content = content.decode(encoding, errors='replace')
            
            config = configparser.ConfigParser()
            config.read_string(ini_content)
            
            # Convert to readable format
            readable_lines = []
            for section in config.sections():
                readable_lines.append(f"[{section}]")
                for key, value in config.items(section):
                    readable_lines.append(f"{key} = {value}")
                readable_lines.append("")
            
            readable_text = "\n".join(readable_lines)
            
            return {
                "text": ini_content,
                "readable_text": readable_text,
                "sections": list(config.sections()),
                "word_count": len(readable_text.split()),
                "document_type": "ini_config_file",
                "encoding": encoding
            }
            
        except Exception as e:
            # Fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_properties(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Java properties files"""
        try:
            encoding = self._detect_encoding(content)
            properties_content = content.decode(encoding, errors='replace')
            
            # Parse properties format
            properties = {}
            readable_lines = []
            
            for line in properties_content.splitlines():
                line = line.strip()
                if line and not line.startswith('#') and '=' in line:
                    key, value = line.split('=', 1)
                    properties[key.strip()] = value.strip()
                    readable_lines.append(f"{key.strip()} = {value.strip()}")
                elif line.startswith('#'):
                    readable_lines.append(line)  # Keep comments
            
            readable_text = "\n".join(readable_lines)
            
            return {
                "text": properties_content,
                "readable_text": readable_text,
                "properties_count": len(properties),
                "word_count": len(readable_text.split()),
                "document_type": "properties_file",
                "encoding": encoding
            }
            
        except Exception as e:
            # Fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_csv(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process CSV/TSV files"""
        try:
            encoding = self._detect_encoding(content)
            csv_content = content.decode(encoding, errors='replace')
            
            # Detect delimiter
            delimiter = '\t' if filename.endswith('.tsv') else ','
            
            # Use csv.Sniffer to detect delimiter if not TSV
            if not filename.endswith('.tsv'):
                try:
                    sniffer = csv.Sniffer()
                    delimiter = sniffer.sniff(csv_content[:1024]).delimiter
                except Exception:
                    delimiter = ','  # fallback to comma
            
            lines = csv_content.splitlines()
            if not lines:
                return {
                    "text": csv_content,
                    "rows": 0,
                    "columns": 0,
                    "document_type": "csv_document",
                    "encoding": encoding
                }
            
            # Use csv.reader for proper parsing
            reader = csv.reader(io.StringIO(csv_content), delimiter=delimiter)
            rows = list(reader)
            
            if not rows:
                return {
                    "text": csv_content,
                    "rows": 0,
                    "columns": 0,
                    "document_type": "csv_document",
                    "encoding": encoding
                }
            
            # Count columns from first row
            columns = len(rows[0]) if rows else 0
            
            # Create readable format
            readable_lines = []
            for i, row in enumerate(rows[:20]):  # Limit to first 20 rows for analysis
                if any(cell.strip() for cell in row):  # Skip empty rows
                    readable_lines.append(f"Row {i+1}: " + " | ".join(row))
            
            if len(rows) > 20:
                readable_lines.append(f"... and {len(rows) - 20} more rows")
            
            readable_text = "\n".join(readable_lines)
            
            return {
                "text": csv_content,
                "readable_text": readable_text,
                "rows": len(rows),
                "columns": columns,
                "delimiter": delimiter,
                "word_count": len(readable_text.split()),
                "document_type": "csv_document",
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process CSV file: {str(e)}")
    
    async def process_document_with_chunking(
        self, 
        file_content: bytes, 
        filename: str, 
        content_type: Optional[str] = None,
        target_provider: str = "groq"
    ) -> Dict[str, Any]:
        """
        Process document with intelligent chunking for AI analysis
        
        Args:
            file_content: Raw file content as bytes
            filename: Original filename
            content_type: MIME content type
            target_provider: Target AI provider for optimization
            
        Returns:
            Processed document with chunked content
        """
        try:
            # 基本的な文書処理
            result = await self.process_document(file_content, filename, content_type)
            
            # プロバイダーに応じたチャンクサイズ調整
            chunk_size = self._get_optimal_chunk_size(target_provider)
            
            # テキストをチャンクに分割
            chunks = self._create_intelligent_chunks(
                result["text"], 
                chunk_size, 
                filename
            )
            
            # 各チャンクにメタデータ付与
            processed_chunks = []
            for i, chunk in enumerate(chunks):
                processed_chunks.append({
                    "chunk_id": i + 1,
                    "content": chunk,
                    "word_count": len(chunk.split()),
                    "char_count": len(chunk),
                    "position": f"{i + 1}/{len(chunks)}"
                })
            
            # 結果に追加
            result.update({
                "chunks": processed_chunks,
                "total_chunks": len(chunks),
                "chunk_size": chunk_size,
                "chunking_method": "intelligent",
                "optimized_for": target_provider
            })
            
            return result
            
        except Exception as e:
            logger.error(f"Failed to process document with chunking: {str(e)}")
            raise DocumentProcessorError(f"Document chunking failed: {str(e)}")
    
    def _get_optimal_chunk_size(self, provider: str) -> int:
        """プロバイダーに応じた最適なチャンクサイズを決定"""
        provider_limits = {
            "groq": 3000,      # Groqは短めに
            "openai": 8000,    # OpenAIは中程度
            "claude": 12000,   # Claudeは長めに
            "gemini": 20000    # Geminiは最大
        }
        return provider_limits.get(provider.lower(), 4000)
    
    def _create_intelligent_chunks(self, text: str, chunk_size: int, filename: str) -> List[str]:
        """
        文書の構造を考慮したインテリジェントなチャンク分割
        """
        chunks = []
        
        # ファイルタイプに応じた分割戦略
        file_extension = Path(filename).suffix.lower()
        
        if file_extension in ['.md', '.txt']:
            # マークダウン・テキストファイルは見出しで分割
            chunks = self._split_by_headings(text, chunk_size)
        elif file_extension in ['.pdf', '.docx']:
            # PDF・Wordは段落で分割
            chunks = self._split_by_paragraphs(text, chunk_size)
        elif file_extension in ['.py', '.js', '.java']:
            # ソースコードは関数・クラス単位で分割
            chunks = self._split_by_code_blocks(text, chunk_size)
        else:
            # その他は汎用的な分割
            chunks = self._split_by_sentences(text, chunk_size)
        
        return chunks
    
    def _split_by_headings(self, text: str, chunk_size: int) -> List[str]:
        """見出し単位での分割"""
        chunks = []
        current_chunk = ""
        
        for line in text.split('\n'):
            # 見出しを検出（#で始まる行）
            if line.strip().startswith('#') and current_chunk:
                if len(current_chunk) > chunk_size:
                    # 大きすぎる場合は文単位で再分割
                    chunks.extend(self._split_by_sentences(current_chunk, chunk_size))
                else:
                    chunks.append(current_chunk.strip())
                current_chunk = line + '\n'
            else:
                current_chunk += line + '\n'
                
                # チャンクサイズ超過時は強制分割
                if len(current_chunk) > chunk_size * 1.5:
                    chunks.extend(self._split_by_sentences(current_chunk, chunk_size))
                    current_chunk = ""
        
        # 最後のチャンクを追加
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _split_by_paragraphs(self, text: str, chunk_size: int) -> List[str]:
        """段落単位での分割"""
        chunks = []
        current_chunk = ""
        
        paragraphs = text.split('\n\n')
        
        for paragraph in paragraphs:
            if len(current_chunk + paragraph) > chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                current_chunk = paragraph + '\n\n'
            else:
                current_chunk += paragraph + '\n\n'
        
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _split_by_code_blocks(self, text: str, chunk_size: int) -> List[str]:
        """コードブロック単位での分割"""
        chunks = []
        current_chunk = ""
        
        lines = text.split('\n')
        current_function = ""
        
        for line in lines:
            # 関数・クラス定義を検出
            if (line.strip().startswith('def ') or 
                line.strip().startswith('class ') or 
                line.strip().startswith('function ')):
                
                if current_chunk and len(current_chunk) > chunk_size:
                    chunks.append(current_chunk.strip())
                    current_chunk = line + '\n'
                else:
                    current_chunk += line + '\n'
            else:
                current_chunk += line + '\n'
                
                if len(current_chunk) > chunk_size * 1.5:
                    chunks.append(current_chunk.strip())
                    current_chunk = ""
        
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _split_by_sentences(self, text: str, chunk_size: int) -> List[str]:
        """文単位での分割（汎用）"""
        chunks = []
        current_chunk = ""
        
        # 日本語の文区切り文字
        sentence_endings = ['。', '！', '？', '.\n', '!\n', '?\n']
        
        sentences = []
        current_sentence = ""
        
        for char in text:
            current_sentence += char
            if any(current_sentence.endswith(ending) for ending in sentence_endings):
                sentences.append(current_sentence.strip())
                current_sentence = ""
        
        if current_sentence.strip():
            sentences.append(current_sentence.strip())
        
        for sentence in sentences:
            if len(current_chunk + sentence) > chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                current_chunk = sentence + ' '
            else:
                current_chunk += sentence + ' '
        
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks

# Global document processor instance
_document_processor: Optional[DocumentProcessor] = None

def get_document_processor() -> DocumentProcessor:
    """Get or create document processor instance"""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor 