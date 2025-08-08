"""
Enhanced Document Processing Service
Phase 9: Comprehensive Knowledge File Support

Supports:
- Microsoft Office: Word (.docx), Excel (.xlsx/.xls), PowerPoint (.pptx)
- PDF documents
- Google Workspace: Docs, Sheets, Slides
- Rich Text Format (.rtf)
- Various text formats: .txt, .log, .xml, .yaml, .ini, .properties, .sql
- Source code files: .py, .js, .html, .css, .java, .cpp, .php, .rb, .go, .rs, .ts
"""

import logging
import io
import re
from typing import Dict, Any, Optional, List, Union
from pathlib import Path
import chardet

# Microsoft Office processing
try:
    from docx import Document as DocxDocument
    import openpyxl
    from pptx import Presentation
    import xlrd
except ImportError:
    DocxDocument = None  # type: ignore
    openpyxl = None  # type: ignore
    Presentation = None  # type: ignore
    xlrd = None  # type: ignore

# PDF processing
try:
    import PyPDF2
    import pdfplumber
except ImportError:
    PyPDF2 = None  # type: ignore
    pdfplumber = None  # type: ignore

# Rich Text and other formats
try:
    from striprtf.striprtf import rtf_to_text
    import yaml  # type: ignore
    from bs4 import BeautifulSoup
    import configparser
except ImportError:
    rtf_to_text = None  # type: ignore
    yaml = None  # type: ignore
    BeautifulSoup = None  # type: ignore
    configparser = None  # type: ignore

# Google API
try:
    from google.oauth2.credentials import Credentials
    from googleapiclient.discovery import build
    from google.auth.transport.requests import Request
    from google_auth_oauthlib.flow import InstalledAppFlow
except ImportError:
    Credentials = None  # type: ignore
    build = None  # type: ignore
    Request = None  # type: ignore
    InstalledAppFlow = None  # type: ignore

logger = logging.getLogger(__name__)

class DocumentProcessorError(Exception):
    """Document processing related errors"""
    pass

class DocumentProcessor:
    """Service for processing various document formats"""
    
    def __init__(self):
        """Initialize document processor"""
        self.supported_formats = {
            '.docx': self._process_docx,
            '.pdf': self._process_pdf,
            '.xlsx': self._process_xlsx,
            '.xls': self._process_xls,
            '.pptx': self._process_pptx,
            '.rtf': self._process_rtf,
            '.txt': self._process_text,
            '.log': self._process_text,
            '.md': self._process_markdown,
            '.markdown': self._process_markdown,
            '.json': self._process_json,
            '.yaml': self._process_yaml,
            '.yml': self._process_yaml,
            '.xml': self._process_xml,
            '.html': self._process_html,
            '.htm': self._process_html,
            '.css': self._process_text,
            '.js': self._process_text,
            '.py': self._process_text,
            '.java': self._process_text,
            '.cpp': self._process_text,
            '.c': self._process_text,
            '.h': self._process_text,
            '.cs': self._process_text,
            '.php': self._process_text,
            '.rb': self._process_text,
            '.go': self._process_text,
            '.rs': self._process_text,
            '.ts': self._process_text,
            '.sh': self._process_text,
            '.sql': self._process_text,
            '.ini': self._process_ini,
            '.properties': self._process_properties,
            '.csv': self._process_csv,
            '.tsv': self._process_csv,
        }
        logger.info("Document processor initialized")
    
    async def process_document(
        self, 
        file_content: bytes, 
        filename: str, 
        content_type: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Process document and extract text content
        
        Args:
            file_content: Raw file content as bytes
            filename: Original filename
            content_type: MIME content type
            
        Returns:
            Processed document information
        """
        try:
            file_path = Path(filename)
            file_extension = file_path.suffix.lower()
            
            if file_extension not in self.supported_formats:
                raise DocumentProcessorError(f"Unsupported file format: {file_extension}")
            
            # Process the document
            processor_func = self.supported_formats[file_extension]
            result = await processor_func(file_content, filename)
            
            # Add metadata
            result.update({
                "filename": filename,
                "file_extension": file_extension,
                "content_type": content_type,
                "file_size": len(file_content),
                "processing_status": "success"
            })
            
            logger.info(f"Successfully processed {filename} ({file_extension})")
            return result
            
        except Exception as e:
            logger.error(f"Failed to process document {filename}: {str(e)}")
            raise DocumentProcessorError(f"Document processing failed: {str(e)}")
    
    def _detect_encoding(self, content: bytes) -> str:
        """Detect text encoding"""
        try:
            detection = chardet.detect(content)
            return detection.get('encoding', 'utf-8') or 'utf-8'
        except:
            return 'utf-8'
    
    async def _process_docx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Word document (.docx)"""
        if DocxDocument is None:
            raise DocumentProcessorError("python-docx not installed")
        
        try:
            doc = DocxDocument(io.BytesIO(content))
            
            # Extract text from paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract text from tables
            tables_content = []
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        tables_content.append(" | ".join(row_text))
            
            full_text = "\n".join(text_content)
            if tables_content:
                full_text += "\n\nTables:\n" + "\n".join(tables_content)
            
            return {
                "text": full_text,
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
                "word_count": len(full_text.split()),
                "document_type": "word_document"
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process Word document: {str(e)}")
    
    async def _process_pdf(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PDF document"""
        text_content = ""
        page_count = 0
        
        # Try pdfplumber first (better for complex layouts)
        if pdfplumber:
            try:
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    page_count = len(pdf.pages)
                    text_parts = []
                    
                    for page in pdf.pages:
                        page_text = page.extract_text()  # type: ignore
                        if page_text:
                            text_parts.append(page_text)
                    
                    text_content = "\n\n".join(text_parts)
                    
                    return {
                        "text": text_content,
                        "pages": page_count,
                        "word_count": len(text_content.split()),
                        "document_type": "pdf_document",
                        "extraction_method": "pdfplumber"
                    }
            except Exception as e:
                logger.warning(f"pdfplumber failed, trying PyPDF2: {str(e)}")
        
        # Fallback to PyPDF2
        if PyPDF2:
            try:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                page_count = len(pdf_reader.pages)
                text_parts = []
                
                for page in pdf_reader.pages:  # type: ignore
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                
                text_content = "\n\n".join(text_parts)
                
                return {
                    "text": text_content,
                    "pages": page_count,
                    "word_count": len(text_content.split()),
                    "document_type": "pdf_document",
                    "extraction_method": "PyPDF2"
                }
            except Exception as e:
                raise DocumentProcessorError(f"Failed to process PDF: {str(e)}")
        
        raise DocumentProcessorError("No PDF processing library available")
    
    async def _process_xlsx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Excel file (.xlsx)"""
        if openpyxl is None:
            raise DocumentProcessorError("openpyxl not installed")
        
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
            sheets_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):  # Skip empty rows
                        sheet_data.append(" | ".join(row_data))
                
                if sheet_data:
                    sheets_content.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_data))
            
            full_text = "\n\n".join(sheets_content)
            
            return {
                "text": full_text,
                "sheets": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames,
                "word_count": len(full_text.split()),
                "document_type": "excel_document"
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process Excel file: {str(e)}")
    
    async def _process_xls(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process legacy Excel file (.xls)"""
        if xlrd is None:
            raise DocumentProcessorError("xlrd not installed")
        
        try:
            workbook = xlrd.open_workbook(file_contents=content)
            sheets_content = []
            
            for sheet_index in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_index)
                sheet_name = workbook.sheet_names()[sheet_index]
                sheet_data = []
                
                for row_idx in range(sheet.nrows):
                    row_data = []
                    for col_idx in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_idx, col_idx)
                        row_data.append(str(cell_value) if cell_value else "")
                    
                    if any(cell.strip() for cell in row_data):
                        sheet_data.append(" | ".join(row_data))
                
                if sheet_data:
                    sheets_content.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_data))
            
            full_text = "\n\n".join(sheets_content)
            
            return {
                "text": full_text,
                "sheets": workbook.nsheets,
                "sheet_names": workbook.sheet_names(),
                "word_count": len(full_text.split()),
                "document_type": "excel_legacy_document"
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process legacy Excel file: {str(e)}")
    
    async def _process_pptx(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process PowerPoint presentation (.pptx)"""
        if Presentation is None:
            raise DocumentProcessorError("python-pptx not installed")
        
        try:
            prs = Presentation(io.BytesIO(content))
            slides_content = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    slides_content.append(f"Slide {slide_idx}:\n" + "\n".join(slide_text))
            
            full_text = "\n\n".join(slides_content)
            
            return {
                "text": full_text,
                "slides": len(prs.slides),
                "word_count": len(full_text.split()),
                "document_type": "powerpoint_document"
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process PowerPoint file: {str(e)}")
    
    async def _process_rtf(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Rich Text Format (.rtf)"""
        if rtf_to_text is None:
            raise DocumentProcessorError("striprtf not installed")
        
        try:
            encoding = self._detect_encoding(content)
            rtf_content = content.decode(encoding)
            text_content = rtf_to_text(rtf_content)
            
            return {
                "text": text_content,
                "word_count": len(text_content.split()),
                "document_type": "rtf_document",
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process RTF file: {str(e)}")
    
    async def _process_text(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process plain text files"""
        try:
            encoding = self._detect_encoding(content)
            text_content = content.decode(encoding)
            
            # Detect file type based on extension
            file_path = Path(filename)
            extension = file_path.suffix.lower()
            
            document_types = {
                '.txt': 'text_document',
                '.log': 'log_file',
                '.py': 'python_source',
                '.js': 'javascript_source',
                '.java': 'java_source',
                '.cpp': 'cpp_source',
                '.c': 'c_source',
                '.h': 'header_file',
                '.cs': 'csharp_source',
                '.php': 'php_source',
                '.rb': 'ruby_source',
                '.go': 'go_source',
                '.rs': 'rust_source',
                '.ts': 'typescript_source',
                '.sh': 'shell_script',
                '.sql': 'sql_script',
                '.css': 'css_stylesheet'
            }
            
            return {
                "text": text_content,
                "lines": len(text_content.splitlines()),
                "word_count": len(text_content.split()),
                "character_count": len(text_content),
                "document_type": document_types.get(extension, 'text_document'),
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process text file: {str(e)}")
    
    async def _process_markdown(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Markdown files"""
        try:
            encoding = self._detect_encoding(content)
            text_content = content.decode(encoding)
            
            # Remove markdown formatting for analysis
            import re
            clean_text = re.sub(r'[#*`\[\]()_~]', '', text_content)
            clean_text = re.sub(r'!\[.*?\]\(.*?\)', '', clean_text)  # Remove images
            clean_text = re.sub(r'\[.*?\]\(.*?\)', '', clean_text)   # Remove links
            
            return {
                "text": text_content,
                "clean_text": clean_text,
                "lines": len(text_content.splitlines()),
                "word_count": len(clean_text.split()),
                "document_type": "markdown_document",
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process Markdown file: {str(e)}")
    
    async def _process_json(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process JSON files"""
        try:
            encoding = self._detect_encoding(content)
            json_content = content.decode(encoding)
            
            # Parse JSON and convert to readable text
            import json
            data = json.loads(json_content)
            
            # Convert JSON to readable format
            readable_text = json.dumps(data, indent=2, ensure_ascii=False)
            
            return {
                "text": json_content,
                "readable_text": readable_text,
                "json_structure": str(type(data).__name__),
                "word_count": len(readable_text.split()),
                "document_type": "json_document",
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process JSON file: {str(e)}")
    
    async def _process_yaml(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process YAML files"""
        if yaml is None:
            return await self._process_text(content, filename)
        
        try:
            encoding = self._detect_encoding(content)
            yaml_content = content.decode(encoding)
            
            # Parse YAML and convert to readable text
            data = yaml.safe_load(yaml_content)
            readable_text = yaml.dump(data, default_flow_style=False, allow_unicode=True)
            
            return {
                "text": yaml_content,
                "readable_text": readable_text,
                "yaml_structure": str(type(data).__name__),
                "word_count": len(readable_text.split()),
                "document_type": "yaml_document",
                "encoding": encoding
            }
            
        except Exception as e:
            # Fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_xml(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process XML files"""
        if BeautifulSoup is None:
            return await self._process_text(content, filename)
        
        try:
            encoding = self._detect_encoding(content)
            xml_content = content.decode(encoding)
            
            # Parse XML and extract text
            soup = BeautifulSoup(xml_content, 'xml')
            text_content = soup.get_text(separator=' ', strip=True)
            
            return {
                "text": xml_content,
                "extracted_text": text_content,
                "word_count": len(text_content.split()),
                "document_type": "xml_document",
                "encoding": encoding
            }
            
        except Exception as e:
            # Fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_html(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process HTML files"""
        if BeautifulSoup is None:
            return await self._process_text(content, filename)
        
        try:
            encoding = self._detect_encoding(content)
            html_content = content.decode(encoding)
            
            # Parse HTML and extract text
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            text_content = soup.get_text(separator=' ', strip=True)
            
            return {
                "text": html_content,
                "extracted_text": text_content,
                "title": soup.title.string if soup.title else None,
                "word_count": len(text_content.split()),
                "document_type": "html_document",
                "encoding": encoding
            }
            
        except Exception as e:
            # Fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_ini(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process INI/config files"""
        if configparser is None:
            return await self._process_text(content, filename)
        
        try:
            encoding = self._detect_encoding(content)
            ini_content = content.decode(encoding)
            
            config = configparser.ConfigParser()
            config.read_string(ini_content)
            
            # Convert to readable format
            readable_lines = []
            for section in config.sections():
                readable_lines.append(f"[{section}]")
                for key, value in config.items(section):
                    readable_lines.append(f"{key} = {value}")
                readable_lines.append("")
            
            readable_text = "\n".join(readable_lines)
            
            return {
                "text": ini_content,
                "readable_text": readable_text,
                "sections": list(config.sections()),
                "word_count": len(readable_text.split()),
                "document_type": "ini_config_file",
                "encoding": encoding
            }
            
        except Exception as e:
            # Fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_properties(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process Java properties files"""
        try:
            encoding = self._detect_encoding(content)
            properties_content = content.decode(encoding)
            
            # Parse properties format
            properties = {}
            readable_lines = []
            
            for line in properties_content.splitlines():
                line = line.strip()
                if line and not line.startswith('#') and '=' in line:
                    key, value = line.split('=', 1)
                    properties[key.strip()] = value.strip()
                    readable_lines.append(f"{key.strip()} = {value.strip()}")
                elif line.startswith('#'):
                    readable_lines.append(line)  # Keep comments
            
            readable_text = "\n".join(readable_lines)
            
            return {
                "text": properties_content,
                "readable_text": readable_text,
                "properties_count": len(properties),
                "word_count": len(readable_text.split()),
                "document_type": "properties_file",
                "encoding": encoding
            }
            
        except Exception as e:
            # Fallback to text processing
            return await self._process_text(content, filename)
    
    async def _process_csv(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Process CSV/TSV files"""
        try:
            encoding = self._detect_encoding(content)
            csv_content = content.decode(encoding)
            
            # Detect delimiter
            delimiter = '\t' if filename.endswith('.tsv') else ','
            
            lines = csv_content.splitlines()
            if not lines:
                return {
                    "text": csv_content,
                    "rows": 0,
                    "columns": 0,
                    "document_type": "csv_document",
                    "encoding": encoding
                }
            
            # Count columns from first line
            first_line = lines[0]
            columns = len(first_line.split(delimiter))
            
            # Create readable format
            readable_lines = []
            for i, line in enumerate(lines[:20]):  # Limit to first 20 rows for analysis
                if line.strip():
                    cells = line.split(delimiter)
                    readable_lines.append(f"Row {i+1}: " + " | ".join(cells))
            
            if len(lines) > 20:
                readable_lines.append(f"... and {len(lines) - 20} more rows")
            
            readable_text = "\n".join(readable_lines)
            
            return {
                "text": csv_content,
                "readable_text": readable_text,
                "rows": len(lines),
                "columns": columns,
                "delimiter": delimiter,
                "word_count": len(readable_text.split()),
                "document_type": "csv_document",
                "encoding": encoding
            }
            
        except Exception as e:
            raise DocumentProcessorError(f"Failed to process CSV file: {str(e)}")

# Global document processor instance
_document_processor: Optional[DocumentProcessor] = None

def get_document_processor() -> DocumentProcessor:
    """Get or create document processor instance"""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor 